"""Generación de archivos de reportería (plan-v3.md, Fases 20-22) —
funciones que reciben datos ya estructurados (título + columnas + filas,
título + secciones, o nodos + aristas, ya armados por Claude) y devuelven
la ruta del archivo escrito en disco, bajo un nombre `<uuid4>.<ext>` bien
determinístico (sin tracking en DB en este v1, ver la nota de
arquitectura del plan — el nombre en sí es la única referencia al
archivo).

Vive acá (no en `rag-mcp/`) por el mismo criterio que
`embeddings.py`/`vector_store.py`: compartido en principio, aunque hoy
solo lo use `chat-rag-mcp` (las tools `report_*` de `rag-mcp/server.py`).
"""

from __future__ import annotations

import csv
import os
import subprocess
import uuid
import xml.etree.ElementTree as ET
from pathlib import Path

import openpyxl

from rag_shared.models import DiagramEdge, DiagramNode, ReportSection

# Coincide con el mount RW nuevo de chat-rag-mcp en compose.yaml
# (`./data/media/reports:/data/media/reports`) — a diferencia del resto de
# `MODEL_CACHE_DIR`/etc., este no lo comparte chat-web via Django settings
# (no hay Django acá): chat-web sirve estos archivos calculando la misma
# ruta como `MEDIA_ROOT / "reports"` (ver apps/core/views.py::download_report).
REPORTS_DIR = Path(os.environ.get("REPORTS_DIR", "/data/media/reports"))

DEFAULT_TITLE = "Reporte"


def _report_path(ext: str) -> Path:
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    return REPORTS_DIR / f"{uuid.uuid4()}.{ext}"


def write_txt(title: str, columns: list[str], rows: list[list[str]]) -> Path:
    """Texto plano legible: título, encabezado, filas separadas por " | "."""
    path = _report_path("txt")
    lines = [title or DEFAULT_TITLE, "", " | ".join(columns), "-" * 40]
    lines += [" | ".join(str(cell) for cell in row) for row in rows]
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return path


def write_csv(title: str, columns: list[str], rows: list[list[str]]) -> Path:
    """CSV estándar: sin línea de título (rompería un import directo a
    Excel/pandas, que esperan que la primera fila sea el encabezado)."""
    path = _report_path("csv")
    with path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(columns)
        writer.writerows(rows)
    return path


def write_xlsx(title: str, columns: list[str], rows: list[list[str]]) -> Path:
    path = _report_path("xlsx")
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    # Los nombres de hoja de Excel tienen un máximo de 31 caracteres —
    # truncar en vez de dejar que openpyxl tire ValueError.
    sheet.title = (title or DEFAULT_TITLE)[:31]
    sheet.append(columns)
    for row in rows:
        sheet.append(row)
    workbook.save(path)
    return path


def _body_lines(body: str) -> list[str]:
    return [line.strip() for line in body.splitlines() if line.strip()]


def write_docx(title: str, sections: list[ReportSection]) -> Path:
    """Documento Word (plan-v3.md, Fase 21): título + una secuencia de
    secciones (encabezado + párrafos/bullets), sin diseño elaborado."""
    import docx

    path = _report_path("docx")
    document = docx.Document()
    document.add_heading(title or DEFAULT_TITLE, level=0)
    for section in sections:
        document.add_heading(section.heading, level=1)
        for line in _body_lines(section.body):
            if line.startswith(("- ", "* ")):
                document.add_paragraph(line[2:], style="List Bullet")
            else:
                document.add_paragraph(line)
    document.save(path)
    return path


def write_pptx(title: str, sections: list[ReportSection]) -> Path:
    """Presentación PowerPoint (plan-v3.md, Fase 21): portada + una
    diapositiva por sección, layout título+contenido de la plantilla
    default de python-pptx (`slide_layouts[1]`, placeholder de contenido
    en el índice 1 — confirmado, no asumido: la plantilla default varía
    entre versiones de python-pptx/PowerPoint)."""
    from pptx import Presentation

    path = _report_path("pptx")
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title or DEFAULT_TITLE

    content_layout = presentation.slide_layouts[1]
    for section in sections:
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = section.heading
        lines = _body_lines(section.body)
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = lines[0] if lines else ""
        for line in lines[1:]:
            text_frame.add_paragraph().text = line

    presentation.save(path)
    return path


# fpdf2 con las fuentes core (Helvetica) solo soporta latin-1 estricto para
# el texto — no cp1252/WinAnsi pese a ser lo que uno esperaría de un PDF:
# cualquier raya larga ("—") o comilla tipográfica ("“"/"”"), MUY comunes en
# la prosa que arma Claude, tira FPDFUnicodeEncodingException y aborta la
# generación entera (confirmado probándolo, no es una suposición). Mapear
# la puntuación "inteligente" más común a su equivalente ASCII/latin-1 antes
# de pasarle nada a fpdf2, y `errors="replace"` como red de seguridad final
# para cualquier otro caracter fuera de rango (ej. emoji) en vez de que
# rompa — se pierde ese caracter puntual, no el reporte entero.
_PDF_CHAR_MAP = str.maketrans(
    {"—": "-", "–": "-", "‘": "'", "’": "'", "“": '"', "”": '"', "…": "...", "•": "-"}
)


def _pdf_safe(text: str) -> str:
    return text.translate(_PDF_CHAR_MAP).encode("latin-1", errors="replace").decode("latin-1")


def write_pdf(title: str, sections: list[ReportSection]) -> Path:
    """PDF (plan-v3.md, Fase 21): título + secuencia de secciones. fpdf2 en
    vez de reportlab — más simple para este caso de uso (texto simple, sin
    layout elaborado en v1)."""
    from fpdf import FPDF, XPos, YPos

    path = _report_path("pdf")
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(0, 10, _pdf_safe(title or DEFAULT_TITLE), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(4)
    for section in sections:
        pdf.set_font("Helvetica", "B", 14)
        pdf.multi_cell(0, 8, _pdf_safe(section.heading), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 11)
        for line in _body_lines(section.body):
            # Los bullets en PDF son un simple "- " a mano (no hay un
            # equivalente directo a "List Bullet" de docx/pptx sin agregar
            # más complejidad de layout) — igual de legible en v1.
            prefix = "- " if line.startswith(("- ", "* ")) else ""
            text = line[2:] if prefix else line
            pdf.multi_cell(0, 6, _pdf_safe(prefix + text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(4)
    pdf.output(str(path))
    return path


# Layout en grilla simple para el .drawio (plan-v3.md, Fase 22) — Claude
# solo manda nodos/aristas, no posiciones. Alcanza: la gracia de un .drawio
# es justo que el usuario lo termina de acomodar a mano en draw.io después,
# a diferencia del .png (Graphviz), que si necesita un layout automático
# real porque no hay "después" — ver write_diagram_png.
_DRAWIO_NODE_W = 160
_DRAWIO_NODE_H = 60
_DRAWIO_GAP = 40
_DRAWIO_COLUMNS = 3


def write_drawio(title: str, nodes: list[DiagramNode], edges: list[DiagramEdge]) -> Path:
    """Diagrama editable (plan-v3.md, Fase 22) — XML `mxGraphModel` armado a
    mano con `xml.etree.ElementTree` (stdlib, sin dependencia nueva). El
    usuario lo abre en diagrams.net/draw.io para seguir editándolo; acá no
    se rasteriza ninguna imagen."""
    path = _report_path("drawio")

    mxfile = ET.Element("mxfile", host="app.diagrams.net")
    diagram = ET.SubElement(mxfile, "diagram", name=(title or DEFAULT_TITLE)[:100], id=str(uuid.uuid4()))
    model = ET.SubElement(
        diagram,
        "mxGraphModel",
        dx="800",
        dy="600",
        grid="1",
        gridSize="10",
        guides="1",
        tooltips="1",
        connect="1",
        arrows="1",
        fold="1",
        page="1",
        pageScale="1",
        pageWidth="850",
        pageHeight="1100",
        math="0",
        shadow="0",
    )
    root = ET.SubElement(model, "root")
    # id "0" y "1" son las dos celdas base que exige el formato drawio (la
    # capa raíz y la primera capa de dibujo) — todo lo demás cuelga de "1".
    ET.SubElement(root, "mxCell", id="0")
    ET.SubElement(root, "mxCell", id="1", parent="0")

    cell_id_by_node_id: dict[str, str] = {}
    for i, node in enumerate(nodes):
        cell_id = f"node-{i}"
        cell_id_by_node_id[node.id] = cell_id
        col, row = i % _DRAWIO_COLUMNS, i // _DRAWIO_COLUMNS
        x = _DRAWIO_GAP + col * (_DRAWIO_NODE_W + _DRAWIO_GAP)
        y = _DRAWIO_GAP + row * (_DRAWIO_NODE_H + _DRAWIO_GAP)
        cell = ET.SubElement(
            root,
            "mxCell",
            id=cell_id,
            value=node.label,
            style="rounded=1;whiteSpace=wrap;html=1;",
            vertex="1",
            parent="1",
        )
        ET.SubElement(
            cell, "mxGeometry", x=str(x), y=str(y), width=str(_DRAWIO_NODE_W), height=str(_DRAWIO_NODE_H), **{"as": "geometry"}
        )

    for i, edge in enumerate(edges):
        source_cell = cell_id_by_node_id.get(edge.source)
        target_cell = cell_id_by_node_id.get(edge.target)
        if source_cell is None or target_cell is None:
            # Referencia a un id de nodo que no vino en `nodes` — se
            # ignora esa arista puntual en vez de romper el diagrama
            # entero (que Claude se haya confundido en un id no debería
            # tirar toda la generación).
            continue
        cell = ET.SubElement(
            root,
            "mxCell",
            id=f"edge-{i}",
            value=edge.label,
            style="edgeStyle=orthogonalEdgeStyle;rounded=0;html=1;",
            edge="1",
            parent="1",
            source=source_cell,
            target=target_cell,
        )
        ET.SubElement(cell, "mxGeometry", relative="1", **{"as": "geometry"})

    tree = ET.ElementTree(mxfile)
    ET.indent(tree, space="  ")
    tree.write(path, encoding="utf-8", xml_declaration=True)
    return path


def _dot_escape(text: str) -> str:
    # Los labels de Graphviz van entre comillas dobles en el .dot — sin
    # escapar, un label con una comilla suelta rompe la sintaxis del
    # archivo (en el mejor caso falla `dot`, en el peor inyecta
    # atributos/nodos que Claude no pidió).
    return text.replace("\\", "\\\\").replace('"', '\\"')


def write_diagram_png(title: str, nodes: list[DiagramNode], edges: list[DiagramEdge]) -> Path:
    """Diagrama rasterizado a PNG vía Graphviz (plan-v3.md, Fase 22) — a
    diferencia de un enfoque tipo Mermaid/drawio-en-vivo (necesitan un
    navegador headless, pesado para este host), Graphviz calcula el layout
    y renderiza sin depender de un browser: `dot` lee el `.dot` armado acá
    por stdin y escribe el PNG directo, vía `subprocess` (necesita el
    binario `dot` instalado — apt `graphviz` en rag-mcp/Dockerfile)."""
    lines = ["digraph G {", "  rankdir=LR;", '  node [shape=box, style="rounded,filled", fillcolor="#f5f5f5", fontname="Helvetica"];']
    if title:
        lines.append(f'  labelloc="t"; fontname="Helvetica"; label="{_dot_escape(title)}";')
    for node in nodes:
        lines.append(f'  "{_dot_escape(node.id)}" [label="{_dot_escape(node.label)}"];')
    for edge in edges:
        label_attr = f' [label="{_dot_escape(edge.label)}"]' if edge.label else ""
        lines.append(f'  "{_dot_escape(edge.source)}" -> "{_dot_escape(edge.target)}"{label_attr};')
    lines.append("}")
    dot_source = "\n".join(lines)

    path = _report_path("png")
    result = subprocess.run(
        ["dot", "-Tpng", "-o", str(path)],
        input=dot_source,
        capture_output=True,
        text=True,
        timeout=30,
    )
    if result.returncode != 0:
        # No debería pasar con un .dot armado por nosotros (no texto crudo
        # de Claude sin escapar) — pero si pasa, mejor un error claro que
        # un archivo PNG vacío/corrupto servido silenciosamente.
        raise RuntimeError(f"graphviz (dot) falló generando el diagrama: {result.stderr.strip()}")
    return path
