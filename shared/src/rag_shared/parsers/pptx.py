from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation

from rag_shared.models import TextBlock


def parse(path: Path) -> Iterator[TextBlock]:
    presentation = Presentation(str(path))

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            yield TextBlock(
                text="\n".join(texts),
                page=slide_number,
                section=f"diapositiva {slide_number}",
            )
